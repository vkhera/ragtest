
import os
import pdfplumber
import chromadb
from chromadb.utils import embedding_functions
import csv
import docx
import pptx

def ingest_pdfs(pdf_dir):
    import shutil
    def get_folder_size(folder):
        total = 0
        for dirpath, dirnames, filenames in os.walk(folder):
            for f in filenames:
                fp = os.path.join(dirpath, f)
                if os.path.isfile(fp):
                    total += os.path.getsize(fp)
        return total

    pending_dir = os.path.join(pdf_dir, "Pending")
    complete_dir = os.path.join(pdf_dir, "Complete")
    os.makedirs(complete_dir, exist_ok=True)

    pending_files = [f for f in os.listdir(pending_dir) if os.path.isfile(os.path.join(pending_dir, f))]
    print(f"Pending files: {pending_files}")
    pending_size = sum(os.path.getsize(os.path.join(pending_dir, f)) for f in pending_files)
    print(f"Total size of files pending ingestion: {pending_size/1024:.2f} KB")

    vectordb_dir = "VectorDB"
    vectordb_size_before = get_folder_size(vectordb_dir)
    print(f"VectorDB size before ingestion: {vectordb_size_before/1024:.2f} KB")

    client = chromadb.PersistentClient(path="VectorDB/ChromaDB")
    ef = embedding_functions.DefaultEmbeddingFunction()
    collection = client.get_or_create_collection("pdf_docs", embedding_function=ef)
    chunk_size = 1000  # characters per chunk
    for fname in pending_files:
        fpath = os.path.join(pending_dir, fname)
        ext = fname.lower().split('.')[-1]
        text = ""
        if fname.lower().endswith('.pdf'):
            try:
                with pdfplumber.open(fpath) as pdf:
                    text = " ".join(page.extract_text() or '' for page in pdf.pages)
            except Exception as e:
                print(f"Error reading PDF {fname}: {e}")
        elif fname.lower().endswith('.txt'):
            try:
                with open(fpath, 'r', encoding='utf-8') as f:
                    text = f.read()
            except Exception as e:
                print(f"Error reading TXT {fname}: {e}")
        elif fname.lower().endswith('.csv'):
            try:
                with open(fpath, 'r', encoding='utf-8') as f:
                    reader = csv.reader(f)
                    text = " ".join([", ".join(row) for row in reader])
            except Exception as e:
                print(f"Error reading CSV {fname}: {e}")
        elif fname.lower().endswith('.docx') or fname.lower().endswith('.doc'):
            try:
                doc = docx.Document(fpath)
                text = " ".join([para.text for para in doc.paragraphs])
            except Exception as e:
                print(f"Error reading DOCX {fname}: {e}")
        elif fname.lower().endswith('.pptx') or fname.lower().endswith('.ppt'):
            try:
                prs = pptx.Presentation(fpath)
                slides_text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slides_text.append(shape.text)
                text = " ".join(slides_text)
            except Exception as e:
                print(f"Error reading PPTX {fname}: {e}")
        if text:
            chunks = [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]
            for idx, chunk in enumerate(chunks):
                chunk_id = f"{fname}_{idx}"
                print(f"Ingesting chunk {chunk_id}...")
                collection.add(documents=[chunk], metadatas=[{"source": fname, "chunk": idx}], ids=[chunk_id])
        # Move file to Complete after ingestion
        shutil.move(fpath, os.path.join(complete_dir, fname))
        print(f"Moved {fname} to Complete folder.")
    vectordb_size_after = get_folder_size(vectordb_dir)
    print(f"VectorDB size after ingestion: {vectordb_size_after/1024:.2f} KB")
    print("Ingestion complete.")
